import os
import io
from pathlib import Path

import streamlit as st
from dotenv import load_dotenv
from openai import OpenAI
from docx import Document
from docx.shared import RGBColor, Pt
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT


load_dotenv()


@st.cache_data(show_spinner=False)
def read_master_reference() -> str:
    """Read the entire reference library from master_reference.md (if present)."""
    path = Path(__file__).with_name("master_reference.md")
    if not path.exists():
        return ""
    content = path.read_text(encoding="utf-8")
    return content.strip()


def _trim_reference_for_prompt(reference_text: str, max_chars: int = 40_000) -> tuple[str, bool]:
    """
    Keep prompts reliable by trimming very large reference libraries to a manageable size.
    Returns (trimmed_text, was_trimmed).
    """
    if len(reference_text) <= max_chars:
        return reference_text, False
    return reference_text[-max_chars:], True


def _detect_prior_relationship(target_group: str) -> bool:
    """Detect if the target group indicates an existing relationship or brand awareness."""
    text = target_group.lower()
    relationship_signals = [
        "already",
        "existing",
        "current customer",
        "in contact",
        "familiar with",
        "know us",
        "working with",
        "using our",
        "past customer",
    ]
    return any(signal in text for signal in relationship_signals)


def _get_methodology_instructions(methodology: str) -> str:
    """
    Return detailed instructions for the chosen sales methodology.
    This becomes the LEADING FRAMEWORK for script generation.
    """
    methodologies = {
        "Standard": """
STANDARD CONVERSATIONAL APPROACH:
- Open with permission-based hook ("Do you have 30 seconds?")
- Build trust through discovery questions
- Tease value proposition without over-explaining
- Use a clear, low-friction CTA
- Balance listening and presenting
- Follow the reference examples' natural flow
""",
        "Poke the Bear": """
POKE THE BEAR METHODOLOGY (Question-Led, Curiosity-Driven):
🎯 CORE PRINCIPLE: Help prospects self-identify pain points through neutral, thought-provoking questions.

STRUCTURE:
1. Hook: Permission + neutral context-setting (NO pitch)
2. Thought-Provoking Questions (3-5): Ask about their current approach, challenges they might not have considered, or industry shifts they're aware of
   - "How are you currently handling X?"
   - "Have you thought about Y?"
   - "What's your take on [emerging trend]?"
3. Knowledge Gap Framing: Introduce the idea that there might be a better way, without pitching yet
4. Soft Value Tease: Only AFTER they've expressed interest or curiosity
5. CTA: Frame as "exploring together" or "seeing if there's a fit"

CRITICAL RULES:
- AVOID early product pitching or feature dumps
- Use questions to SPARK CURIOSITY, not interrogate
- Let them talk 70% of the time
- Focus on THEIR world, not your solution
- Lower defenses by being consultative, not salesy
- Make them WANT to learn more before you offer
""",
        "Pitch the Product": """
PITCH THE PRODUCT METHODOLOGY (Solution-Centric, Direct Value):
🎯 CORE PRINCIPLE: Lead with features, benefits, and proof points to demonstrate clear superiority.

STRUCTURE:
1. Hook: Permission + immediate value framing ("We help companies like yours solve X")
2. Why Now: Industry challenge + our solution's relevance
3. Feature-Benefit Stack: 3-4 specific capabilities with tangible outcomes
   - "Our [feature] means you get [benefit]"
   - Use data, case studies, proof points
4. Competitive Differentiation: What makes this better than alternatives
5. Strong CTA: Confidence-driven ask ("Let's get you set up with a demo")

CRITICAL RULES:
- LEAD with value and capabilities early (don't wait)
- Use PROOF POINTS (data, case studies, customer wins)
- Be direct about superiority and differentiation
- Focus on solving INDUSTRY-LEVEL challenges, not just individual pain
- Speak with authority and confidence
- Assume they have the pain; position the solution immediately
"""
    }
    return methodologies.get(methodology, methodologies["Standard"]).strip()


def _extract_text_from_pdf(file) -> str:
    """Extract text from a PDF file using pdfplumber."""
    try:
        import pdfplumber
        with pdfplumber.open(file) as pdf:
            text = ""
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n\n"
            return text.strip()
    except Exception as e:
        return f"[Error extracting PDF: {e}]"


def _extract_text_from_docx(file) -> str:
    """Extract text from a DOCX file using python-docx."""
    try:
        from docx import Document
        doc = Document(file)
        text = ""
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text += paragraph.text + "\n"
        # Also extract text from tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        text += cell.text + " "
                text += "\n"
        return text.strip()
    except Exception as e:
        return f"[Error extracting DOCX: {e}]"


def _extract_text_from_pptx(file) -> str:
    """Extract text from a PPTX file using python-pptx."""
    try:
        from pptx import Presentation
        prs = Presentation(file)
        text = ""
        for slide_num, slide in enumerate(prs.slides, start=1):
            text += f"[Slide {slide_num}]\n"
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text += shape.text + "\n"
            text += "\n"
        return text.strip()
    except Exception as e:
        return f"[Error extracting PPTX: {e}]"


def extract_text_from_uploaded_files(uploaded_files) -> str:
    """
    Extract text from uploaded files (PDF, DOCX, PPTX).
    Returns combined text from all files.
    """
    if not uploaded_files:
        return ""
    
    combined_text = ""
    for uploaded_file in uploaded_files:
        file_name = uploaded_file.name
        file_extension = file_name.lower().split(".")[-1]
        
        combined_text += f"\n\n{'='*60}\n"
        combined_text += f"SOURCE: {file_name}\n"
        combined_text += f"{'='*60}\n\n"
        
        if file_extension == "pdf":
            combined_text += _extract_text_from_pdf(uploaded_file)
        elif file_extension == "docx":
            combined_text += _extract_text_from_docx(uploaded_file)
        elif file_extension == "pptx":
            combined_text += _extract_text_from_pptx(uploaded_file)
        else:
            combined_text += f"[Unsupported file type: {file_extension}]"
        
        combined_text += "\n\n"
    
    return combined_text.strip()


def scrape_website_content(url: str) -> tuple[str, str]:
    """
    Scrape main text content from a website URL, filtering out navigation and UI elements.
    Returns (content, status_message) tuple.
    Status can be "success", "error", or "empty".
    """
    if not url or not url.strip():
        return "", "empty"
    
    url = url.strip()
    # Add https:// if no protocol specified
    if not url.startswith(("http://", "https://")):
        url = "https://" + url
    
    try:
        import requests
        from bs4 import BeautifulSoup
        
        # Set a reasonable timeout and user agent
        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36"
        }
        response = requests.get(url, headers=headers, timeout=10)
        response.raise_for_status()
        
        soup = BeautifulSoup(response.content, "html.parser")
        
        # Aggressively remove non-content elements
        for element in soup([
            "script", "style", "nav", "footer", "header", 
            "aside", "form", "button", "iframe", "noscript",
            "meta", "link"
        ]):
            element.decompose()
        
        # Remove common UI element classes/IDs
        ui_patterns = [
            "nav", "menu", "sidebar", "footer", "header", "cookie",
            "login", "signin", "signup", "cart", "search", "breadcrumb",
            "advertisement", "banner", "popup", "modal"
        ]
        
        for pattern in ui_patterns:
            for elem in soup.find_all(attrs={"class": lambda x: x and pattern in str(x).lower()}):
                elem.decompose()
            for elem in soup.find_all(attrs={"id": lambda x: x and pattern in str(x).lower()}):
                elem.decompose()
        
        # Get text content
        text = soup.get_text(separator="\n", strip=True)
        
        # Filter out common navigation/button phrases (Swedish and English)
        noise_phrases = [
            "logga in", "log in", "sign in", "kontakta oss", "contact us",
            "bli kund", "become a customer", "läs mer", "read more",
            "klicka här", "click here", "meny", "menu", "sök", "search",
            "cookies", "gdpr", "integritetspolicy", "privacy policy",
            "hem", "home", "om oss", "about us", "nyheter", "news"
        ]
        
        # Clean up and filter lines
        lines = []
        for line in text.splitlines():
            line_stripped = line.strip()
            
            # Skip empty lines
            if not line_stripped:
                continue
            
            # Skip very short lines (likely navigation)
            if len(line_stripped) < 15:
                continue
            
            # Skip lines that are just noise phrases
            if line_stripped.lower() in noise_phrases:
                continue
            
            # Skip lines with only numbers or special characters
            if not any(c.isalpha() for c in line_stripped):
                continue
            
            lines.append(line_stripped)
        
        clean_text = "\n".join(lines)
        
        if not clean_text or len(clean_text) < 100:
            return "", "error: Insufficient content extracted from website"
        
        # Limit to reasonable size (first 8,000 chars for cleaner summarization)
        if len(clean_text) > 8000:
            clean_text = clean_text[:8000]
        
        return clean_text, "success"
    
    except requests.exceptions.Timeout:
        return "", "error: Website request timed out (>10s)"
    except requests.exceptions.HTTPError as e:
        return "", f"error: HTTP {e.response.status_code} - {e.response.reason}"
    except requests.exceptions.ConnectionError:
        return "", "error: Could not connect to website"
    except requests.exceptions.RequestException as e:
        return "", f"error: {str(e)}"
    except Exception as e:
        return "", f"error: Unexpected error - {str(e)}"


def summarize_website_content(website_text: str, website_url: str) -> str:
    """
    Use OpenAI to summarize website content into clean, actionable insights.
    Returns a concise summary of core business value, services, and USP.
    """
    if not website_text or not website_text.strip():
        return ""
    
    try:
        api_key = os.getenv("OPENAI_API_KEY")
        if not api_key:
            return website_text[:500]  # Fallback to truncated raw text
        
        client = OpenAI(api_key=api_key)
        
        prompt = f"""Analyze this website content and provide a CONCISE summary (max 150 words) covering:
1. Core business value proposition
2. Main services/products offered
3. Key unique selling points (USP)
4. Target audience (if mentioned)

CRITICAL: Only include business-relevant insights. Do NOT include:
- Navigation menu items
- UI button text (like "Logga in", "Kontakta oss", "Bli kund")
- Footer information
- Random website words

Website URL: {website_url}

Website Content:
{website_text[:4000]}

Provide a clean, professional summary in 2-3 short paragraphs."""

        response = client.chat.completions.create(
            model="gpt-4o-mini",
            temperature=0.3,
            max_tokens=300,
            messages=[
                {
                    "role": "system",
                    "content": "You extract and summarize key business information from website text. Be concise and focus only on value propositions, services, and unique selling points."
                },
                {
                    "role": "user",
                    "content": prompt
                }
            ]
        )
        
        summary = response.choices[0].message.content.strip()
        return summary if summary else website_text[:500]
    
    except Exception as e:
        # Fallback to truncated raw text if summarization fails
        return website_text[:500] + "..."


def export_to_word(
    guide_content: str,
    client_info: str = "",
) -> tuple[bytes, str]:
    """
    Export the generated guide to a Word document using template.docx.
    Returns (docx_bytes, status_message).
    
    THREE-PLACEHOLDER MAPPING:
    - {{CONTENT}} → Main Guide (Inledning, Discovery, Värdeförslag, Avslutning / Boka möte, Invändningar)
    - {{TALKING_POINTS}} → Detailed Talking Points (with page break before)
    - {{CLIENT_INFO}} → Website summary
    
    FORMATTING:
    - Headers: BLACK and BOLD (not green)
    - Bullets: Standard • character
    - Spacing: Single line, compact layout (2-3 pages max)
    """
    template_path = Path(__file__).with_name("template.docx")
    
    try:
        # Check if template exists and is loadable
        if not template_path.exists() or not template_path.is_file():
            # Create a basic document from scratch
            doc = Document()
            
            # Manual heading (avoids style errors)
            p = doc.add_paragraph()
            run = p.add_run("Discussion Guide")
            run.font.bold = True
            run.font.size = Pt(16)
            run.font.color.rgb = RGBColor(0, 0, 0)
            
            if client_info:
                p = doc.add_paragraph()
                run = p.add_run("Klientinformation")
                run.font.bold = True
                run.font.size = Pt(14)
                run.font.color.rgb = RGBColor(0, 0, 0)
                doc.add_paragraph(client_info)
            
            p = doc.add_paragraph()
            run = p.add_run("Guide Content")
            run.font.bold = True
            run.font.size = Pt(14)
            run.font.color.rgb = RGBColor(0, 0, 0)
            _add_markdown_to_doc(doc, guide_content)
            
            status = "warning: Mall hittades ej - använder standardformatering."
        else:
            # Attempt to load template
            try:
                doc = Document(template_path)
                
                # ========================================
                # CLEAN DOCUMENT: Accept all tracked changes and disable tracking
                # This prevents strikethrough text when replacing placeholders
                # ========================================
                try:
                    # Accept all existing tracked changes (removes strikethroughs)
                    if hasattr(doc, '_element'):
                        # Remove all delete (w:del) elements
                        for del_elem in doc._element.xpath('//w:del'):
                            del_elem.getparent().remove(del_elem)
                        
                        # Accept all insert (w:ins) elements by unwrapping them
                        for ins_elem in doc._element.xpath('//w:ins'):
                            parent = ins_elem.getparent()
                            for child in ins_elem:
                                parent.insert(parent.index(ins_elem), child)
                            parent.remove(ins_elem)
                        
                        # Turn off track changes for future edits
                        track_changes_elements = doc._element.xpath('//w:trackRevisions')
                        for tc in track_changes_elements:
                            tc.getparent().remove(tc)
                except:
                    # If we can't clean track changes, continue anyway
                    pass
                
            except Exception as load_error:
                # If template is corrupted, create new document
                doc = Document()
                
                # Manual heading (avoids style errors)
                p = doc.add_paragraph()
                run = p.add_run("Discussion Guide")
                run.font.bold = True
                run.font.size = Pt(16)
                run.font.color.rgb = RGBColor(0, 0, 0)
                
                if client_info:
                    p = doc.add_paragraph()
                    run = p.add_run("Klientinformation")
                    run.font.bold = True
                    run.font.size = Pt(14)
                    run.font.color.rgb = RGBColor(0, 0, 0)
                    doc.add_paragraph(client_info)
                
                p = doc.add_paragraph()
                run = p.add_run("Guide Content")
                run.font.bold = True
                run.font.size = Pt(14)
                run.font.color.rgb = RGBColor(0, 0, 0)
                _add_markdown_to_doc(doc, guide_content)
                status = f"warning: Mall kunde inte läsas ({str(load_error)[:50]}) - använder standardformatering."
                buffer = io.BytesIO()
                doc.save(buffer)
                buffer.seek(0)
                return buffer.getvalue(), status
            
            # ========================================
            # TEMPLATE PLACEHOLDER MAPPING
            # ========================================
            # Template Order:
            # 1. {{CONTENT}} → Main Discussion Guide (Hook, Discovery, Value Prop, Objections)
            # 2. {{CLIENT_INFO}} → Website summary (middle section)
            # 3. {{TALKING_POINTS}} → Detailed Talking Points (very end)
            # ========================================
            
            # Step 1: Split generated content into Main Guide and Talking Points
            main_guide, talking_points = _split_guide_and_talking_points(guide_content)
            
            # ========================================
            # STEP-BY-STEP PLACEHOLDER REPLACEMENT
            # Process in order: CONTENT → CLIENT_INFO → TALKING_POINTS
            # ========================================
            
            # Step 2: Replace {{CONTENT}} with Main Discussion Guide
            # (Hook, Discovery, Value Prop, Objections - NO Talking Points)
            content_replaced = False
            for i, paragraph in enumerate(doc.paragraphs):
                if "{{CONTENT}}" in paragraph.text:
                    # Insert content BEFORE placeholder, then remove placeholder
                    _replace_content_in_place(doc, main_guide, i)
                    # Remove the placeholder paragraph after insertion
                    p_element = paragraph._element
                    p_element.getparent().remove(p_element)
                    content_replaced = True
                    break
            
            if not content_replaced:
                # Fallback: append at end if {{CONTENT}} not found
                doc.add_page_break()
                p = doc.add_paragraph()
                run = p.add_run("Discussion Guide")
                run.font.bold = True
                run.font.size = Pt(14)
                run.font.color.rgb = RGBColor(0, 0, 0)
                _add_markdown_to_doc(doc, main_guide)
            
            # Step 3: Replace {{CLIENT_INFO}} with website summary
            # Clean replacement to avoid track changes strikethroughs
            for paragraph in doc.paragraphs:
                if "{{CLIENT_INFO}}" in paragraph.text:
                    # Clean the client_info first
                    clean_info = client_info.strip() if client_info else ""
                    
                    # Clear the paragraph completely and rebuild it
                    paragraph.clear()
                    if clean_info:
                        run = paragraph.add_run(clean_info)
                        run.font.name = 'Calibri'  # Match reference file
                        run.font.size = Pt(12)
                        run.font.color.rgb = RGBColor(0, 0, 0)
                    
                    break  # Only replace first occurrence
            
            # Step 4: Replace {{TALKING_POINTS}} with detailed Talking Points
            # IMPORTANT: Must search fresh after previous insertions
            talking_points_replaced = False
            if talking_points:
                # Re-enumerate paragraphs (indices have shifted after {{CONTENT}} insertion)
                for i, paragraph in enumerate(doc.paragraphs):
                    if "{{TALKING_POINTS}}" in paragraph.text:
                        # NO page break - let template handle spacing
                        
                        # Insert talking points BEFORE placeholder, then remove placeholder
                        _replace_content_in_place(doc, talking_points, i)
                        # Remove the placeholder paragraph after insertion
                        p_element = paragraph._element
                        p_element.getparent().remove(p_element)
                        talking_points_replaced = True
                        break
                
                # Fallback: append at end if {{TALKING_POINTS}} not found
                if not talking_points_replaced:
                    doc.add_page_break()
                    _add_markdown_to_doc(doc, talking_points)
            
            status = "success"
        
        # Save to bytes buffer
        buffer = io.BytesIO()
        doc.save(buffer)
        buffer.seek(0)
        
        return buffer.getvalue(), status
    
    except Exception as e:
        # Last resort: create minimal document
        try:
            doc = Document()
            
            # Manual heading (avoids style errors)
            p = doc.add_paragraph()
            run = p.add_run("Discussion Guide")
            run.font.bold = True
            run.font.size = Pt(16)
            run.font.color.rgb = RGBColor(0, 0, 0)
            
            doc.add_paragraph(guide_content)
            buffer = io.BytesIO()
            doc.save(buffer)
            buffer.seek(0)
            return buffer.getvalue(), f"warning: Export with limited formatting - {str(e)[:50]}"
        except:
            return b"", f"error: {str(e)}"


def _clean_generated_content(content: str) -> str:
    """
    Remove any AI-generated text that should NOT be in the output.
    Specifically targets sections that are already in the template.
    
    This is a SAFETY FILTER that catches content the AI shouldn't have generated.
    """
    import re
    
    # Strip leading/trailing whitespace
    content = content.strip()
    
    # Pattern to detect and remove "If not now, when?" sections (entire section)
    # Matches various formats: "If not now, when?", "Om inte nu, när?", etc.
    forbidden_patterns = [
        r'##?\s*(?:If not now|Om inte nu)[^\n]*\n[^#]*',  # Section header + content
        r'(?:If not now, when\?|Om inte nu, när\?)[^\n]*\n[^#]*',  # Direct text
    ]
    
    for pattern in forbidden_patterns:
        content = re.sub(pattern, '', content, flags=re.IGNORECASE | re.MULTILINE)
    
    # Remove standalone phrases that might appear
    forbidden_phrases = [
        "if not now, when?",
        "om inte nu, när?",
        "lead status",
        "projektstatus",
        "very important!",
        "mycket viktigt!"
    ]
    
    lines = content.split("\n")
    cleaned_lines = []
    
    for line in lines:
        line_lower = line.strip().lower()
        # Skip lines that contain forbidden phrases
        if not any(phrase in line_lower for phrase in forbidden_phrases):
            cleaned_lines.append(line)
    
    return "\n".join(cleaned_lines).strip()


def _split_guide_and_talking_points(guide_content: str) -> tuple[str, str]:
    """
    Split the generated guide into two parts:
    1. Main Guide (Inledning, Discovery, Värdeförslag, Invändningar, Avslutning / Boka möte)
    2. Talking Points section
    
    Filters out any "If not now, when?", "Closing", or "Avslut" sections that shouldn't be generated.
    
    Returns (main_guide, talking_points)
    """
    # First, clean the content to remove any forbidden text
    guide_content = _clean_generated_content(guide_content)
    
    lines = guide_content.split("\n")
    main_guide_lines = []
    talking_points_lines = []
    
    in_talking_points = False
    skip_section = False
    
    for line in lines:
        line_stripped = line.strip().lower()
        
        # Detect sections that should be filtered out (already in template)
        if line.startswith("## "):
            # Check for forbidden sections
            forbidden_keywords = [
                "if not now", "om inte nu", "closing",
                "uppföljning", "follow", "next steps", "nästa steg",
                "lead status", "projektstatus"
            ]
            
            if any(keyword in line_stripped for keyword in forbidden_keywords):
                skip_section = True
                continue  # Skip this header and its content
            
            # Detect Talking Points section (allowed)
            if "talking" in line_stripped or "punkter" in line_stripped or "samtalspunkter" in line_stripped:
                in_talking_points = True
                skip_section = False
                continue  # Skip the "## Talking Points" header itself (template has it)
            else:
                # Reset skip flag for allowed sections
                skip_section = False
        
        # Skip lines in forbidden sections
        if skip_section:
            continue
        
        # Add to appropriate section
        if in_talking_points:
            talking_points_lines.append(line)
        else:
            main_guide_lines.append(line)
    
    main_guide = "\n".join(main_guide_lines).strip()
    talking_points = "\n".join(talking_points_lines).strip()
    
    return main_guide, talking_points


def _replace_content_in_place(doc: Document, markdown_text: str, position_index: int):
    """
    Replace placeholder with content AT THE EXACT POSITION in the template.
    Matches the reference file formatting EXACTLY:
    - Headers: Calibri 14pt Bold, space_before=8pt (first=0pt), space_after=6pt
    - Horizontal lines: Gray underscores, 8pt, space_before=6pt, space_after=6pt
    - Body text: Calibri 12pt, space_before=1pt, space_after=1pt
    - Bullets: Standard bullets (•) with 12pt text
    
    Args:
        doc: The Word Document object
        markdown_text: The markdown content to insert
        position_index: The paragraph index where the placeholder is located
    
    Strategy:
    - Insert all content paragraphs BEFORE the paragraph at position_index
    - This pushes the placeholder paragraph down
    - Caller removes the placeholder after insertion is complete
    """
    from docx.shared import Pt, RGBColor
    from docx.enum.text import WD_LINE_SPACING
    import re
    
    # Strip leading/trailing whitespace and newlines to prevent blank pages
    markdown_text = markdown_text.strip()
    
    lines = markdown_text.split("\n")
    previous_was_header = False
    previous_was_bullet = False
    previous_was_numbered = False
    in_talking_points = False
    in_objections = False  # Track if we're in Invändningar (Objections) section
    first_header = True  # Track if this is the very first header
    previous_was_bold_objection = False  # Track if previous line was a bold objection statement
    
    # Get reference to the placeholder paragraph
    # We'll insert all content before it
    placeholder_paragraph = doc.paragraphs[position_index] if position_index < len(doc.paragraphs) else None
    
    first_header = True  # Track if this is the very first header
    content_started = False  # Track if we've started adding content
    
    for line in lines:
        line_stripped = line.strip()
        
        # Skip completely empty lines, especially at the start (remove top margin)
        if not line_stripped:
            # Don't add empty paragraphs at the start of the document
            if not content_started:
                continue
            
            if previous_was_header:
                if placeholder_paragraph:
                    paragraph = placeholder_paragraph.insert_paragraph_before()
                else:
                    paragraph = doc.add_paragraph()
                paragraph.paragraph_format.space_after = Pt(3)
                previous_was_header = False
            continue
        
        # Check for horizontal rule - EXACT match to reference file
        if line_stripped.startswith("---"):
            if placeholder_paragraph:
                paragraph = placeholder_paragraph.insert_paragraph_before()
            else:
                paragraph = doc.add_paragraph()
            # Horizontal line formatting from reference: 8pt gray text, 6pt spacing before/after
            paragraph.paragraph_format.space_before = Pt(6)
            paragraph.paragraph_format.space_after = Pt(6)
            paragraph.paragraph_format.line_spacing_rule = WD_LINE_SPACING.SINGLE
            run = paragraph.add_run("_" * 40)  # 40 underscores as in reference
            run.font.size = Pt(8)
            run.font.color.rgb = RGBColor(128, 128, 128)  # Gray color
            previous_was_bullet = False
            previous_was_bold_objection = False
            continue
        
        # Check for H2 headers (##)
        if line_stripped.startswith("## "):
            header_text = line_stripped[3:].strip()
            
            content_started = True  # Mark that content has started
            
            # MANUAL HEADER FORMATTING (bypasses 'Heading 2' style error)
            if placeholder_paragraph:
                paragraph = placeholder_paragraph.insert_paragraph_before()
            else:
                paragraph = doc.add_paragraph()
            
            # Track which section we're in
            if "talking" in header_text.lower() or "punkter" in header_text.lower():
                in_talking_points = True
                in_objections = False
            elif "invändning" in header_text.lower() or "objection" in header_text.lower():
                in_objections = True
                in_talking_points = False
            else:
                in_talking_points = False
                in_objections = False
            
            # Header formatting - EXACT match to reference file
            if first_header:
                paragraph.paragraph_format.space_before = Pt(0)  # No top margin for first header
                first_header = False
            else:
                paragraph.paragraph_format.space_before = Pt(8)  # 8pt space before non-first headers
            
            paragraph.paragraph_format.space_after = Pt(6)  # 6pt space after headers (reference spec)
            paragraph.paragraph_format.line_spacing_rule = WD_LINE_SPACING.SINGLE
            
            # Add the header text: Calibri 14pt Bold Black (EXACT match to reference)
            run = paragraph.add_run(header_text)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0, 0, 0)  # Black
            run.font.size = Pt(14)
            run.font.name = 'Calibri'
            
            previous_was_header = True
            previous_was_bullet = False
            previous_was_numbered = False
            previous_was_bold_objection = False
            continue
        
        # Check for numbered points (1., 2., 3., etc.)
        numbered_match = re.match(r'^(\d+)\.\s+(.+)$', line_stripped)
        if numbered_match:
            number = numbered_match.group(1)
            text = numbered_match.group(2)
            
            content_started = True  # Mark that content has started
            
            # Add the numbered point with proper spacing
            if placeholder_paragraph:
                paragraph = placeholder_paragraph.insert_paragraph_before()
            else:
                paragraph = doc.add_paragraph()
            
            # Talking Points spacing: 1 blank line BEFORE number (after previous bullets)
            # and 1 blank line AFTER number (before this section's bullets)
            if in_talking_points:
                if previous_was_numbered or previous_was_bullet:
                    paragraph.paragraph_format.space_before = Pt(12)  # 1 blank line before number
                else:
                    paragraph.paragraph_format.space_before = Pt(0)  # First item: no space before
                paragraph.paragraph_format.space_after = Pt(12)  # 1 blank line after number
            else:
                paragraph.paragraph_format.space_before = Pt(0)
                paragraph.paragraph_format.space_after = Pt(1)
            paragraph.paragraph_format.line_spacing_rule = WD_LINE_SPACING.SINGLE
            
            number_run = paragraph.add_run(f"{number}. ")
            number_run.font.name = 'Calibri'
            number_run.font.size = Pt(12)
            number_run.font.bold = True
            number_run.font.color.rgb = RGBColor(0, 0, 0)
            
            _add_formatted_text_to_paragraph(paragraph, text, font_size=12)
            
            previous_was_numbered = True
            previous_was_bullet = False
            previous_was_header = False
            previous_was_bold_objection = False
            continue
        
        # Check for bullets (-, *, >)
        if line_stripped.startswith(("- ", "* ", "> ")):
            bullet_text = line_stripped[2:].strip()
            
            content_started = True  # Mark that content has started
            
            if placeholder_paragraph:
                paragraph = placeholder_paragraph.insert_paragraph_before()
            else:
                paragraph = doc.add_paragraph()
            
            # COMPACT SPACING for bullets
            if in_talking_points and previous_was_numbered:
                paragraph.paragraph_format.space_before = Pt(2)
            else:
                paragraph.paragraph_format.space_before = Pt(1)
            
            paragraph.paragraph_format.space_after = Pt(1)
            paragraph.paragraph_format.line_spacing_rule = WD_LINE_SPACING.SINGLE
            paragraph.paragraph_format.left_indent = Pt(20)
            
            # Use simple bullet character - match reference formatting
            bullet_run = paragraph.add_run("• ")
            bullet_run.font.name = 'Calibri'
            bullet_run.font.size = Pt(12)
            bullet_run.font.color.rgb = RGBColor(0, 0, 0)
            
            _add_formatted_text_to_paragraph(paragraph, bullet_text, font_size=12)
            
            previous_was_bullet = True
            previous_was_numbered = False
            previous_was_header = False
            previous_was_bold_objection = False
            continue
        
        # Regular paragraph text
        content_started = True  # Mark that content has started
        
        # Check if this is a bold objection statement (e.g., **"Vi har redan en lösning"**)
        is_bold_objection = (in_objections and 
                            line_stripped.startswith("**") and 
                            line_stripped.endswith("**") and
                            len(line_stripped) > 4)
        
        # Insert blank paragraph BEFORE bold objection in Objections section (for visual separation)
        if is_bold_objection and previous_was_bold_objection:
            if placeholder_paragraph:
                blank_para = placeholder_paragraph.insert_paragraph_before()
            else:
                blank_para = doc.add_paragraph()
            blank_para.paragraph_format.space_before = Pt(0)
            blank_para.paragraph_format.space_after = Pt(0)
            blank_para.paragraph_format.line_spacing_rule = WD_LINE_SPACING.SINGLE
        
        if placeholder_paragraph:
            paragraph = placeholder_paragraph.insert_paragraph_before()
        else:
            paragraph = doc.add_paragraph()
        
        paragraph.paragraph_format.space_before = Pt(1)
        paragraph.paragraph_format.space_after = Pt(1)
        paragraph.paragraph_format.line_spacing_rule = WD_LINE_SPACING.SINGLE
        
        _add_formatted_text_to_paragraph(paragraph, line_stripped, font_size=12)
        
        previous_was_bullet = False
        previous_was_numbered = False
        previous_was_header = False
        previous_was_bold_objection = is_bold_objection  # Track for next iteration


def _add_markdown_to_doc(doc: Document, markdown_text: str):
    """
    Convert markdown-style text to Word document formatting.
    Matches reference file EXACTLY:
    - Headers: Calibri 14pt Bold, space_before=8pt (first=0pt), space_after=6pt
    - Horizontal lines: Gray underscores, 8pt, space_before=6pt, space_after=6pt
    - Body text: Calibri 12pt, space_before=1pt, space_after=1pt
    - Bullets: Standard bullets (•) with 12pt Calibri text
    """
    from docx.shared import Pt
    from docx.enum.text import WD_LINE_SPACING
    import re
    
    # Strip leading/trailing whitespace and newlines to prevent blank pages
    markdown_text = markdown_text.strip()
    
    lines = markdown_text.split("\n")
    previous_was_header = False
    previous_was_bullet = False
    previous_was_numbered = False
    in_talking_points = False
    in_objections = False  # Track if we're in Invändningar (Objections) section
    first_header = True  # Track if this is the very first header
    content_started = False  # Track if we've started adding content
    previous_was_bold_objection = False  # Track if previous line was a bold objection statement
    
    for line in lines:
        line_stripped = line.strip()
        
        # Skip completely empty lines, especially at the start (remove top margin)
        if not line_stripped:
            # Don't add empty paragraphs at the start of the document
            if not content_started:
                continue
            
            # Only add spacing after headers or between major sections
            if previous_was_header:
                paragraph = doc.add_paragraph()
                paragraph.paragraph_format.space_after = Pt(3)
                previous_was_header = False
            continue
        
        # Check for horizontal rule - EXACT match to reference file
        if line_stripped.startswith("---"):
            paragraph = doc.add_paragraph()
            # Horizontal line formatting from reference: 8pt gray text, 6pt spacing before/after
            paragraph.paragraph_format.space_before = Pt(6)
            paragraph.paragraph_format.space_after = Pt(6)
            paragraph.paragraph_format.line_spacing_rule = WD_LINE_SPACING.SINGLE
            run = paragraph.add_run("_" * 40)  # 40 underscores as in reference
            run.font.size = Pt(8)
            run.font.color.rgb = RGBColor(128, 128, 128)  # Gray color
            previous_was_bullet = False
            previous_was_bold_objection = False
            continue
        
        # Check for H2 headers (##)
        if line_stripped.startswith("## "):
            header_text = line_stripped[3:].strip()
            
            content_started = True  # Mark that content has started
            
            # MANUAL HEADER FORMATTING (bypasses 'Heading 2' style error)
            # Create a regular paragraph and format it as a header manually
            paragraph = doc.add_paragraph()
            
            # Track which section we're in
            if "talking" in header_text.lower() or "punkter" in header_text.lower():
                in_talking_points = True
                in_objections = False
            elif "invändning" in header_text.lower() or "objection" in header_text.lower():
                in_objections = True
                in_talking_points = False
            else:
                in_talking_points = False
                in_objections = False
            
            # Header formatting - EXACT match to reference file
            if first_header:
                paragraph.paragraph_format.space_before = Pt(0)  # No top margin for first header
                first_header = False
            else:
                paragraph.paragraph_format.space_before = Pt(8)  # 8pt space before non-first headers
            
            paragraph.paragraph_format.space_after = Pt(6)  # 6pt space after headers (reference spec)
            paragraph.paragraph_format.line_spacing_rule = WD_LINE_SPACING.SINGLE
            
            # Add the header text: Calibri 14pt Bold Black (EXACT match to reference)
            run = paragraph.add_run(header_text)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0, 0, 0)  # Black
            run.font.size = Pt(14)
            run.font.name = 'Calibri'
            
            previous_was_header = True
            previous_was_bullet = False
            previous_was_numbered = False
            previous_was_bold_objection = False
            continue
        
        # Check for numbered points (1., 2., 3., etc.) - especially in Talking Points
        numbered_match = re.match(r'^(\d+)\.\s+(.+)$', line_stripped)
        if numbered_match:
            number = numbered_match.group(1)
            text = numbered_match.group(2)
            
            # Add the numbered point with proper spacing
            paragraph = doc.add_paragraph()
            
            # Talking Points spacing: 1 blank line BEFORE number (after previous bullets)
            # and 1 blank line AFTER number (before this section's bullets)
            if in_talking_points:
                if previous_was_numbered or previous_was_bullet:
                    paragraph.paragraph_format.space_before = Pt(12)  # 1 blank line before number
                else:
                    paragraph.paragraph_format.space_before = Pt(0)  # First item: no space before
                paragraph.paragraph_format.space_after = Pt(12)  # 1 blank line after number
            else:
                paragraph.paragraph_format.space_before = Pt(0)
                paragraph.paragraph_format.space_after = Pt(1)
            paragraph.paragraph_format.line_spacing_rule = WD_LINE_SPACING.SINGLE
            
            # Format as: "1. Text"
            number_run = paragraph.add_run(f"{number}. ")
            number_run.font.name = 'Calibri'
            number_run.font.size = Pt(12)
            number_run.font.bold = True
            number_run.font.color.rgb = RGBColor(0, 0, 0)
            
            _add_formatted_text_to_paragraph(paragraph, text, font_size=12)
            
            previous_was_numbered = True
            previous_was_bullet = False
            previous_was_header = False
            previous_was_bold_objection = False
            continue
        
        # Check for bullets (-, *, >)
        if line_stripped.startswith(("- ", "* ", "> ")):
            bullet_text = line_stripped[2:].strip()
            
            # Create compact paragraph with bullet
            paragraph = doc.add_paragraph()
            
            # COMPACT SPACING for bullets (tight grouping under numbered points)
            paragraph.paragraph_format.space_before = Pt(0)
            paragraph.paragraph_format.space_after = Pt(1)
            paragraph.paragraph_format.line_spacing_rule = WD_LINE_SPACING.SINGLE
            
            # Detect indent level
            indent_level = (len(line) - len(line.lstrip())) // 2
            
            # Add bullet character based on indent
            if indent_level == 0:
                bullet_char = "• "
            elif indent_level == 1:
                bullet_char = "  ◦ "
            else:
                bullet_char = "    ▪ "
            
            # Add bullet character - match reference formatting
            bullet_run = paragraph.add_run(bullet_char)
            bullet_run.font.name = 'Calibri'
            bullet_run.font.size = Pt(12)
            bullet_run.font.color.rgb = RGBColor(0, 0, 0)
            
            # Add formatted text after bullet
            _add_formatted_text_to_paragraph(paragraph, bullet_text, font_size=12)
            
            previous_was_bullet = True
            previous_was_header = False
            # Don't reset previous_was_numbered - bullets can be under numbered points
            continue
        
        # Regular paragraph - compact formatting with SINGLE SPACING
        content_started = True  # Mark that content has started
        
        # Check if this is a bold objection statement (e.g., **"Vi har redan en lösning"**)
        is_bold_objection = (in_objections and 
                            line_stripped.startswith("**") and 
                            line_stripped.endswith("**") and
                            len(line_stripped) > 4)
        
        # Insert blank paragraph BEFORE bold objection in Objections section (for visual separation)
        if is_bold_objection and previous_was_bold_objection:
            blank_para = doc.add_paragraph()
            blank_para.paragraph_format.space_before = Pt(0)
            blank_para.paragraph_format.space_after = Pt(0)
            blank_para.paragraph_format.line_spacing_rule = WD_LINE_SPACING.SINGLE
        
        paragraph = doc.add_paragraph()
        paragraph.paragraph_format.space_before = Pt(1)
        paragraph.paragraph_format.space_after = Pt(1)
        paragraph.paragraph_format.line_spacing_rule = WD_LINE_SPACING.SINGLE
        
        _add_formatted_text_to_paragraph(paragraph, line_stripped, font_size=12)
        
        previous_was_bold_objection = is_bold_objection  # Track for next iteration
        
        previous_was_bullet = False
        previous_was_header = False
        previous_was_numbered = False


def _add_formatted_text_to_paragraph(paragraph, text: str, font_size: int = 12):
    """
    Add text to paragraph with bold formatting for **text** markers.
    Matches reference file: Calibri 12pt, with bold for emphasized terms.
    """
    from docx.shared import Pt, RGBColor
    
    parts = text.split("**")
    
    for i, part in enumerate(parts):
        if not part:
            continue
        
        run = paragraph.add_run(part)
        run.font.name = 'Calibri'  # Match reference file font
        run.font.size = Pt(font_size)
        run.font.color.rgb = RGBColor(0, 0, 0)  # Black text
        
        # Every odd index (1, 3, 5...) should be bold
        if i % 2 == 1:
            run.font.bold = True


def _build_prompt(
    product: str,
    goal: str,
    target_group: str,
    personas: str,
    tone_of_voice: str,
    reference_examples: str,
    audience_mode: str,
    feedback: str,
    additional_reference: str,
    source_materials: str,
    website_context: str,
    sales_methodology: str,
    do_not_talk_about: str,
    strict_qualifying_questions: str,
) -> str:
    audience_instruction = (
        "New Customer: lead with trust-building, highlight the status-quo problem, and why change now."
        if audience_mode == "New Customer"
        else "Upsell: lean on the existing relationship, show how this is the natural next step for their size/revenue, and reflect established familiarity."
    )
    
    # Detect prior relationship
    has_prior_relationship = _detect_prior_relationship(target_group)
    relationship_instruction = ""
    if has_prior_relationship:
        relationship_instruction = (
            "\n⚠️ CRITICAL: The Target Group indicates a prior relationship or brand awareness. "
            "SKIP all brand introductions and 'what is [company]' explanations. "
            "Start the conversation from relevance, not from zero. "
            "Assume they know who you are—focus on WHY THIS MATTERS NOW for their specific context."
        )
    
    # Elevate feedback and additional notes to hard constraints
    hard_constraints = []
    if feedback.strip():
        hard_constraints.append(f"Recent call feedback (HIGH PRIORITY CONSTRAINT): {feedback.strip()}")
    if additional_reference.strip():
        hard_constraints.append(f"Additional constraint notes (MANDATORY): {additional_reference.strip()}")
    
    constraints_block = ""
    if hard_constraints:
        constraints_block = "\n\n🚨 MANDATORY CONSTRAINTS (failure to follow = script rejected):\n" + "\n".join(
            f"- {constraint}" for constraint in hard_constraints
        )
    
    rewrite_objection = ""
    feedback_lower = feedback.lower()
    if any(keyword in feedback_lower for keyword in ["busy", "time", "no time", "too busy", "calendar"]):
        rewrite_objection = (
            "\n- Objection Handling MUST directly address 'I don't have time' for this persona. "
            "Make it concise, respectful of time, and offer a frictionless next step."
        )
    
    # Get methodology instructions
    methodology_instructions = _get_methodology_instructions(sales_methodology)
    
    # Build critical negative constraints block
    negative_constraints_block = ""
    if do_not_talk_about and do_not_talk_about.strip():
        negative_constraints_block = f"""
🚨 CRITICAL NEGATIVE CONSTRAINTS (HIGHEST PRIORITY):
You are STRICTLY FORBIDDEN from mentioning, discussing, or referencing the following topics:
{do_not_talk_about.strip()}

⛔ ABSOLUTE RULES:
- Even if these topics appear in Source Materials or Client Website, DO NOT include them in the script
- If a topic is in the "Do Not Talk About" list, COMPLETELY AVOID IT regardless of any other instructions
- This constraint OVERRIDES all other guidance including reference examples and methodology
- Violation of this constraint = script rejected

"""
    
    # Build mandatory qualification block
    qualification_block = ""
    if strict_qualifying_questions and strict_qualifying_questions.strip():
        qualification_block = f"""
⚠️ MANDATORY QUALIFICATION CRITERIA (HIGH PRIORITY):
The user has specified non-negotiable criteria that MUST be verified during discovery:
{strict_qualifying_questions.strip()}

✅ REQUIRED ACTIONS:
- Include specific questions in the Discovery section to verify these criteria
- Make these questions feel natural but remain firm on the requirements
- Visually mark these as mandatory qualification by wrapping them in a callout box:
  
  > 📋 **KVALIFICERING (Mandatory):**
  > - [Question to verify criterion 1]
  > - [Question to verify criterion 2]

- Place this KVALIFICERING box within or immediately after the Discovery Questions section
- The salesperson must know these are non-negotiable checkpoints

"""
    
    return f"""
You are crafting a COLD OUTREACH calling script for a sales campaign.

THIS IS NOT MARKETING COPY. THIS IS NOT A WEBSITE. THIS IS A REAL PHONE CALL.
Be punchy, direct, and conversational. Eliminate copywriter fluff and corporate bullshit.

{negative_constraints_block}{qualification_block}
🎯 SALES METHODOLOGY (THE LEADING FRAMEWORK FOR THIS SCRIPT):
{methodology_instructions}

⚠️ CRITICAL: The chosen Sales Methodology is your PRIMARY guide for structure and tone.
- If "Poke the Bear" is selected: AVOID early pitching. Focus on questions and curiosity. Even with increased length, maintain the question-led approach.
- If "Pitch the Product" is selected: LEAD with value and capabilities upfront. Use the additional length to add proof points and detailed features.
- If "Standard" is selected: Follow the balanced, conversational reference approach. Scale up naturally with more dialogue and detail.
- 📏 The comprehensive length requirement does NOT override your methodology — it ENHANCES it with more depth while maintaining the core framework.

Tone of Voice:
{tone_of_voice.strip() or 'Professional yet conversational. Confident, concise, and helpful.'}

📖 REFERENCE EXAMPLES — THE GOLD STANDARD (master_reference.md):
{reference_examples.strip() or 'N/A'}

⭐ CRITICAL INSTRUCTION FOR REFERENCE EXAMPLES:
- These are your GOLD STANDARD for length, depth, structure, and comprehensiveness
- MATCH the word count and detail level of these examples (typically 400-600+ words)
- If references include specific sub-sections, bullets, or detailed scenarios, MIRROR that structure
- Pay attention to how they balance conversational script with comprehensive Talking Points
- Notice the pacing: multiple questions, explanations, and natural dialogue flow
- The "Talking points" section in references is MANDATORY in your output — replicate this format with 5-8 numbered points and sub-bullets
- Do NOT generate shorter, summarized versions — match the FULL depth of the references

Deep Knowledge Base (uploaded source materials — use to make the pitch more specific, accurate, and compelling):
{source_materials.strip() if source_materials else 'N/A — rely on Product field and your knowledge'}

📚 INSTRUCTION: If source materials are provided above, consult them to find:
- Specific technical specs, features, or capabilities
- Unique value propositions and differentiators
- Pain points and solutions mentioned in the documents
- Data, case studies, or proof points that align with the Goal
Use these details to make the script more credible and tailored, but keep it conversational (not a data dump).

Client Brand & Service Context (scraped from client's website):
{website_context.strip() if website_context else 'N/A — no website provided'}

🌐 INSTRUCTION: If website context is provided above, use it to ensure:
- The pitch's tone and brand voice match how the client presents themselves
- Service descriptions are 100% accurate to the client's own messaging
- You reference their actual offerings, not assumed ones
- The script feels like it came from someone who genuinely understands their brand
This is a complement to other source materials — cross-reference for accuracy.

Context Inputs (STRICTLY FOLLOW THESE):
- Product: {product.strip() or 'N/A'}
- Goal: {goal}
- Target Group: {target_group.strip() or 'N/A'}
- Personas: {personas.strip() or 'N/A'}
- Audience Focus: {audience_mode} ({audience_instruction})
{relationship_instruction}
{constraints_block}

Output Structure (comprehensive cold call script matching reference depth):

📋 FORMATTING & STYLE REQUIREMENTS (MANDATORY for conversational, glanceable guides):

🎨 CONVERSATIONAL FLOW (Anti-Script):
- AVOID robotic, perfect prose — use conversational markers: "Hördu...", "Egentligen...", "Kortfattat..."
- Focus on FLOW and FEEL, not exact wording
- Make it sound like natural human conversation, not reading from a script
- Use short, punchy phrases in bullets (not long sentences)
- Each point = standalone thought with LOTS of whitespace

✨ GLANCEABLE FORMATTING:
1. Use ## (H2 headers) for EVERY major section
2. Insert --- (horizontal rule) BETWEEN every section
3. Use **bold keywords** WITHIN sentences so caller sees core message at a glance
4. Bullet-led structure (-) for everything (short phrases, not paragraphs)
5. For KEY MOMENTS (Hook, Ask), provide 2-3 ALTERNATIVES formatted as:
   - **Alt A (Casual):** [short phrase]
   - **Alt B (Direct):** [short phrase]
   - **Alt C (Question-led):** [short phrase]
6. Avoid text blocks — maximize whitespace between ideas

🚨 INVÄNDNINGAR (OBJECTIONS) - CRITICAL FORMATTING:
- Format objections with **bold customer quote** followed by 2-3 clean response bullets
- Example format:
  **"Vi har redan en lösning"**
  - Response option 1...
  - Response option 2...
- DO NOT include meta-instructions like _(Listen, then)_, _(Or acknowledge)_, [short timeframe], etc.
- Keep responses CLEAN, NATURAL, and ready-to-use without editorial notes

📋 OUTPUT STRUCTURE - CRITICAL:
- START DIRECTLY with the first ## header (e.g., ## Inledning)
- DO NOT add introductory paragraphs or concluding sentences
- DO NOT write "Med denna guide får ni..." or "This guide will help you..." at the end

🚨 SECTIONS TO GENERATE (EXACTLY THESE, NO MORE):
1. ## Inledning / Hook (with alternatives)
2. ## Discovery / Kvalificering (questions)
3. ## Värdeförslag (value proposition)
4. ## Invändningar (objections with clean responses)
5. ## Avslutning / Boka möte (meeting booking with 2-3 soft alternatives)
6. ## Talking Points (numbered detailed points)

⛔ DO NOT GENERATE THESE (already in template):
- "If not now, when?" section
- Lead status definitions (A-lead, B-lead, C-lead)
- STOP immediately after ## Talking Points

EXAMPLE FORMATTING STRUCTURE:

## Inledning / Hook

**Alt A (Casual):** Hördu, har du **2 minuter** att prata om [topic]?

**Alt B (Direct):** Ringer från [company] — handlar om **[value prop]**. Passar det?

**Alt C (Permission):** Stör jag? Kan jag få låna **30 sekunder**?

---

## Är du bekant med oss?

- **Ja** → Perfekt! Då vet du att vi jobbar med [X]
- **Nej** → Inget konstigt — vi bygger **[solution]** för [target group]

---

## Vad vi gör

- Egentligen ganska enkelt: **[core offering]**
- Vi tar hand om **allt** från [A] till [B]
- Fokus på **[key benefit]** — inte bara ännu en leverantör

---

## Discovery Frågor

**Får jag fråga:**

- Hur jobbar ni med **[area]** idag?
- Har ni funderat på **[challenge/opportunity]**?
- Vad använder ni för **[solution type]** nu?
- Planerar ni att **[future direction]**?

_(Pick 2-3 that fit the flow — not all)_

---

## Värdeförslaget

**Kortfattat:**

- Vi erbjuder **"[Solution]-as-a-Service"**
  - Ni slipper **investering** → istället **prenumeration**
  - Vi fixar **installation, drift, underhåll**
  - **Förutsägbara** kostnader — ingen överraskningar

- Prismodell: **månadsavgift** istället för att ni binder kapital

---

## Uppföljning

**Beroende på svar, fråga:**

- Har ni funderat på att nyttja **[alternative approach]**?
- Skulle det underlätta om ni kunde **[specific benefit]**?

---

## Avslut / Mötesfrågan

**Alt A (Soft):** Kan du tänka dig ha **30 minuter** under [timeframe] för att vi kör igenom det närmare?

**Alt B (Direct):** Bokar vi in **en halvtimme** nästa vecka?

**Alt C (Exploratory):** Låter det relevant? Kan vi få **15 minuter** för att se om det finns en fit?

---

## Invändningar

**"Vi har redan en lösning":**
- Toppen! Hur funkar den för er idag?
- _(Listen, then)_ Kan vara värt att jämföra — ofta hittar vi **[specific benefit]** som ni missar

**"Skicka info på mail":**
- Absolut, men tar bara **2 minuter** att förklara — sen vet du om det är relevant alls
- _(Or acknowledge)_ Gör jag, men vill du att jag bokar in **10 minuter** samtidigt?

**"Har inte tid":**
- Förstår — därför jag frågar om **[short timeframe]** senare, inte nu
- Går det bättre om vi siktar på **[specific future time]**?

---

## Avslutning / Boka möte

**Alt A (Soft):** Kan du tänka dig ha **30 minuter** under [timeframe] för att vi kör igenom det närmare?

**Alt B (Direct):** Skulle du vara öppen att låna mig **en halvtimme** kommande veckor?

**Alt C (Conversational):** Hur låter det om vi bokar in **30 minuter** framöver över Teams?

---

## Talking Points

1. **Prenumeration istället för investering**
   - Företag slipper **stor initial kostnad**
   - Vi tar **hela kontraktet** = mindre risk för er
   - **Förutsägbara** månadsavgifter

2. **Helhetslösning — allt ingår**
   - Installation, **drift**, underhåll
   - **Bemannad driftcentral** 24/7
   - Tekniker på plats inom **4 timmar**

3. **[Continue with 5-8 numbered points]**
   - Each with 2-4 sub-bullets
   - Bold the **key differentiators**

CONTENT REQUIREMENTS (bullet-led, conversational style):

- **Hook / Permission:** 2-3 alternatives (Casual / Direct / Question-led) — short, punchy
- **Brand Awareness Check:** Bullet format with Yes/No paths
- **What We Do:** 3-4 short bullets with **bold keywords** — not paragraphs
- **Discovery Questions:** 3-5 bullets, each a standalone question. Add note: "_(Pick 2-3 that fit)_"
- **Value Proposition:** Bullet structure with **"Kortfattat:"** intro, then sub-bullets with bold terms
- **Additional Questions:** 1-2 bullets starting with "Beroende på svar..." or "Om de säger X..."
- **Objection Handling:** 3-4 common objections with conversational responses (use "Toppen!", "Förstår...", "Absolut...")
- **Avslutning / Boka möte:** 2-3 soft alternatives for the meeting ask using phrases like "Kan du tänka dig...", "Skulle du vara öppen...", "Hur låter det..."
- **Talking Points:** 5-8 numbered with **bold headings** and 2-4 sub-bullets each. Bold all **key differentiators**

📏 LENGTH & DEPTH REQUIREMENTS (comprehensive BUT glanceable):
- Match the DEPTH of master_reference.md (400-600 words MINIMUM) but in bullet/short-phrase format
- Use **bullets and short phrases**, NOT long sentences or paragraphs (exception: brief context within bullets)
- Each bullet = standalone thought with whitespace around it
- Provide **multiple alternatives** for Hook and Close sections (2-3 each)
- Talking Points must be comprehensive (150-200 words) with 5-8 numbered items and sub-bullets
- More bullets with **bold keywords** > fewer long paragraphs
- MAXIMIZE whitespace — avoid text blocks, make it airy and scannable

CRITICAL CONSTRAINTS:

🗣️ CONVERSATIONAL FLOW (ANTI-SCRIPT):
- AVOID "reading from a script" feel — make it airy, glanceable, human
- Use conversational markers: "Hördu...", "Egentligen...", "Kortfattat...", "Toppen!", "Förstår..."
- Focus on FLOW, not perfect wording — give variations and alternatives
- Sound like a **human conversation**, not corporate prose
- Short phrases in bullets > long sentences

✨ FORMATTING (MANDATORY):
- Bullet-led structure for EVERYTHING (except Talking Points headings)
- **Bold keywords** within bullets for glanceability
- Provide 2-3 **alternatives** for Hook and Close sections
- Use ## headers, --- separators between sections
- MAXIMIZE whitespace — each bullet is a standalone thought
- The output should be easy to scan mid-call, not read word-for-word

🎯 CONTENT:
- This is COLD OUTREACH, not a brochure. Be direct, punchy, human.
- STRICTLY follow the Target Group, Personas, and any Additional Reference Notes.
- If the Target Group indicates prior relationship/brand awareness, START FROM RELEVANCE, not introductions.
- Be SPECIFIC to the inputs above. Generic scripts will be rejected.
- Match the reference pacing and "borrow X minutes" style.
- USE MASTER_REFERENCE.MD AS THE GOLD STANDARD FOR DEPTH (but reformat into bullets)
- Even with increased length, content MUST reflect the chosen Sales Methodology
{rewrite_objection}
""".strip()


def _audience_mode(personas: str) -> str:
    text = personas.lower()
    upsell_keywords = ["upsell", "expand", "expansion", "existing", "customer base", "current customer"]
    if any(keyword in text for keyword in upsell_keywords):
        return "Upsell"
    return "New Customer"


def generate_guide(
    product: str,
    goal: str,
    target_group: str,
    personas: str,
    tone_of_voice: str,
    additional_reference: str,
    feedback: str,
    source_materials: str = "",
    website_context: str = "",
    sales_methodology: str = "Standard",
    do_not_talk_about: str = "",
    strict_qualifying_questions: str = "",
) -> str:
    master_reference = read_master_reference()
    trimmed_master_reference, _ = _trim_reference_for_prompt(master_reference)
    reference_examples = trimmed_master_reference
    if additional_reference.strip():
        reference_examples = (
            (
                trimmed_master_reference + "\n\n---\n\nAdditional reference notes:\n" + additional_reference
            ).strip()
            if trimmed_master_reference
            else additional_reference.strip()
        )

    audience_mode = _audience_mode(personas)
    prompt = _build_prompt(
        product=product,
        goal=goal,
        target_group=target_group,
        personas=personas,
        tone_of_voice=tone_of_voice,
        reference_examples=reference_examples,
        audience_mode=audience_mode,
        feedback=feedback,
        additional_reference=additional_reference,
        source_materials=source_materials,
        website_context=website_context,
        sales_methodology=sales_methodology,
        do_not_talk_about=do_not_talk_about,
        strict_qualifying_questions=strict_qualifying_questions,
    )

    try:
        api_key = os.getenv("OPENAI_API_KEY")
        if not api_key:
            # UI should handle this before calling generation, but keep a safe guard here.
            raise ValueError("Missing OPENAI_API_KEY")

        client = OpenAI(api_key=api_key)
        response = client.chat.completions.create(
            model="gpt-4o-mini",
            temperature=0.7,
            messages=[
                {
                    "role": "system",
                    "content": (
                        "You create sharp, CONVERSATIONAL, GLANCEABLE cold calling guides. "
                        "You are NOT writing marketing copy or rigid scripts—this is a REAL, flowing human conversation. "
                        "\n\n"
                        "🚨 CRITICAL NEGATIVE CONSTRAINTS (ABSOLUTE HIGHEST PRIORITY): "
                        "If the user provides a 'Do Not Talk About' list, you are STRICTLY FORBIDDEN from mentioning those topics. "
                        "Even if they appear in source materials, website content, or reference examples, COMPLETELY AVOID THEM. "
                        "This constraint OVERRIDES ALL other instructions. Violation = script rejected. "
                        "\n\n"
                        "⚠️ MANDATORY QUALIFICATION: "
                        "If the user specifies 'Strict Qualifying Questions', you MUST include these in the Discovery section. "
                        "Wrap them in a callout box with '> 📋 **KVALIFICERING (Mandatory):**' so the salesperson knows they are non-negotiable checkpoints. "
                        "Make the questions feel natural but firm on requirements. "
                        "\n\n"
                        "🗣️ ANTI-SCRIPT APPROACH (CRITICAL): "
                        "AVOID the 'reading from a script' feel. Make it airy, bullet-led, scannable. "
                        "Use conversational markers like 'Hördu...', 'Egentligen...', 'Kortfattat...', 'Toppen!', 'Förstår...'. "
                        "Focus on FLOW and FEEL, not perfect wording. Provide ALTERNATIVES for key moments (Hook, Close). "
                        "Use **bold keywords** within bullets so caller can glance and see core message instantly. "
                        "\n\n"
                        "📏 LENGTH & STRUCTURE: Match the DEPTH of master_reference.md BUT keep it COMPACT for Word export (2-3 pages max). "
                        "Use bullets and short phrases, NOT long sentences or paragraphs. Each bullet = standalone thought. "
                        "Group Alt A/B/C options CLOSELY together (single section, not spread out). Keep 'Talking Points' punchy and concise (5-8 items). "
                        "IMPORTANT: Minimize excessive whitespace between sections — the guide should be dense but scannable. "
                        "\n\n"
                        "🎯 SALES METHODOLOGY: The chosen framework is your PRIMARY guide. "
                        "If 'Poke the Bear' → question-led, no early pitch. If 'Pitch the Product' → lead with value/proof. "
                        "If 'Standard' → balanced conversational. The bullet format ENHANCES the methodology, not replaces it. "
                        "\n\n"
                        "✨ FORMATTING (MANDATORY): "
                        "## headers for sections, --- separators, **bold keywords** within bullets, bullet-led structure for everything. "
                        "Provide alternatives: '**Alt A (Casual):** [phrase]', '**Alt B (Direct):** [phrase]'. "
                        "Make it easy to scan mid-call—not read word-for-word. "
                        "\n\n"
                        "🚨 INVÄNDNINGAR (OBJECTIONS) FORMATTING: "
                        "When writing the objections section, use CLEAN dialogue only. "
                        "Format: **Bold the customer's objection** and list 2-3 clean response options as bullets below. "
                        "DO NOT include meta-instructions like _(Listen, then)_, _(Or acknowledge)_, [short timeframe], etc. "
                        "Keep it natural and ready-to-use without editorial notes. "
                        "\n\n"
                        "📅 AVSLUTNING / BOKA MÖTE (CLOSING / MEETING BOOKING): "
                        "This is the CTA section where you ask for the meeting. CRITICAL formatting: "
                        "Provide 2-3 alternatives using conversational Swedish phrases like 'Kan du tänka dig...', 'Skulle du vara öppen...', 'Hur låter det...'. "
                        "Request 30 minutes ('30 minuter' or 'en halvtimme'). "
                        "Mention timeframe: 'under [timeframe]', 'kommande veckor', 'framöver'. "
                        "Optionally mention format: 'över Teams', 'för att presentera närmare'. "
                        "Make it soft, non-pushy, professional. This is separate from 'If not now, when?' which is already in the template. "
                        "\n\n"
                        "📋 OUTPUT STRUCTURE: "
                        "START DIRECTLY with the first ## header (e.g., ## Inledning). "
                        "GENERATE EXACTLY THESE SECTIONS IN ORDER: ## Inledning, ## Discovery, ## Värdeförslag, ## Avslutning / Boka möte, ## Invändningar, ## Talking Points. "
                        "DO NOT generate 'If not now, when?' or lead status sections - these are already in the template. "
                        "STOP after ## Talking Points. No introductory or concluding text. "
                        "\n\n"
                        "Analyze reference examples for depth/tone/pacing, but REFORMAT into bullets with bold keywords. "
                        "STRICTLY follow Target Group, Personas, Additional Notes, Feedback. "
                        "If prior relationship indicated, SKIP introductions. Be specific, human, conversational—not corporate."
                    ),
                },
                {"role": "user", "content": prompt},
            ],
        )
        return (response.choices[0].message.content or "").strip()
    except Exception as exc:  # pragma: no cover - surfaced to UI
        return f"Error generating guide: {exc}"


def render_generator_view() -> None:
    """Render the main generator interface."""
    st.title("📞 Discussion Guide Generator")
    st.caption("Generate ready-to-use, conversational cold calling scripts tailored to your campaign.")

    with st.sidebar:
        st.header("⚙️ Campaign Setup")
        
        # === GROUP 1: BASICS ===
        st.subheader("📋 Basics")
        
        product = st.text_area(
            "Product / Offering",
            height=150,
            placeholder="Describe the product, service, or solution you're pitching...",
            help="Be specific about what you're selling and key benefits.",
            key="input_product",
            value=st.session_state.input_product
        )
        
        goal = st.selectbox(
            "Campaign Goal",
            options=["Leads", "Meetings", "Workshops"],
            index=["Leads", "Meetings", "Workshops"].index(st.session_state.input_goal),
            help="What's the primary objective of this outreach?",
            key="input_goal"
        )
        
        target_group = st.text_area(
            "Target Group",
            height=150,
            placeholder="Industries, company size, revenue, employee count, geographic focus...",
            help="Define who you're targeting with as much detail as possible.",
            key="input_target_group",
            value=st.session_state.input_target_group
        )
        
        personas = st.text_area(
            "Personas / Decision Makers",
            height=150,
            placeholder="Job titles, roles, responsibilities, new customers vs. upsell...",
            help="Who will you be speaking with? Include seniority level and buying authority.",
            key="input_personas",
            value=st.session_state.input_personas
        )
        
        strict_qualifying_questions = st.text_area(
            "Strict Qualifying Questions ⚠️",
            height=120,
            placeholder='E.g., "Must have a turnover of at least 50M SEK", "Must use system X today"...',
            help="🎯 Non-negotiable criteria that MUST be verified during discovery.",
            key="input_strict_qualifying_questions",
            value=st.session_state.input_strict_qualifying_questions
        )
        
        st.divider()
        
        # === GROUP 2: STRATEGY ===
        st.subheader("🎯 Strategy")
        
        client_website_url = st.text_input(
            "Client Website URL (optional)",
            placeholder="https://example.com",
            help="Scrape the website to match client's brand voice and services.",
            key="input_client_website_url",
            value=st.session_state.input_client_website_url
        )
        
        sales_methodology = st.selectbox(
            "Sales Methodology",
            options=["Standard", "Poke the Bear", "Pitch the Product"],
            index=["Standard", "Poke the Bear", "Pitch the Product"].index(st.session_state.input_sales_methodology),
            help="Choose the strategic framework that shapes how the script is structured and delivered.",
            key="input_sales_methodology"
        )
        
        # Dynamic helper text based on methodology
        methodology_descriptions = {
            "Standard": "💬 **Conversational & Balanced** — Respectful permission-based opening, discovery questions, value tease, and clear CTA.",
            "Poke the Bear": "🤔 **Question-Led & Thought-Provoking** — Use neutral questions to help prospects self-identify pain points and knowledge gaps.",
            "Pitch the Product": "🚀 **Solution-Centric & Direct** — Lead with features, benefits, and USPs early. Use proof points to demonstrate superiority.",
        }
        st.info(methodology_descriptions[sales_methodology])
        
        tone_of_voice = st.text_input(
            "Tone of Voice",
            value=st.session_state.input_tone_of_voice,
            help="Describe the desired communication style.",
            key="input_tone_of_voice"
        )
        
        st.divider()
        
        # === GROUP 3: GUARDRAILS ===
        st.subheader("🚨 Guardrails")
        
        do_not_talk_about = st.text_area(
            "Do Not Talk About (Negative Constraints)",
            height=120,
            placeholder='E.g., "Don\'t mention we are Swedish", "Avoid pricing", "Do not mention competitor X"...',
            help="🚨 Anything written here will be STRICTLY AVOIDED by the AI—even if it appears in source materials.",
            key="input_do_not_talk_about",
            value=st.session_state.input_do_not_talk_about
        )
        
        additional_reference = st.text_area(
            "Additional Constraints & Notes",
            height=120,
            placeholder="E.g., 'Avoid corporate jargon', 'They already know our brand', 'Skip introductions'...",
            help="These will be treated as MANDATORY constraints by the AI.",
            key="input_additional_reference",
            value=st.session_state.input_additional_reference
        )
        
        feedback = st.text_area(
            "Recent Call Feedback",
            height=120,
            placeholder="E.g., 'They keep saying they are too busy' or other objections heard live.",
            help="High-priority feedback that will adjust the script generation.",
            key="input_feedback",
            value=st.session_state.input_feedback
        )
        
        st.divider()
        
        # === SOURCE MATERIALS ===
        with st.expander("📁 Source Materials (Optional)", expanded=False):
            uploaded_files = st.file_uploader(
                "Upload Materials (PPTX, PDF, DOCX)",
                type=["pdf", "docx", "pptx"],
                accept_multiple_files=True,
                help="Product decks, case studies, or technical docs to enrich the pitch.",
            )
            
            # Reference file status
            master_reference = read_master_reference()
            if master_reference:
                trimmed_reference, was_trimmed = _trim_reference_for_prompt(master_reference)
                st.success(f"✓ Loaded reference library ({len(master_reference):,} chars)")
                if was_trimmed:
                    st.caption(f"Using last {len(trimmed_reference):,} characters")
            else:
                st.warning("⚠️ No reference file found")
        
        # Generate Button (always visible at bottom)
        st.divider()
        generate_clicked = st.button(
            "🚀 Generate Script",
            use_container_width=True,
            type="primary"
        )

    st.subheader("Cold Calling Script")
    if generate_clicked:
        if not os.getenv("OPENAI_API_KEY"):
            st.error(
                "OPENAI_API_KEY is missing. Add it to your `.env` file (or your environment) and restart the app."
            )
            return

        # Scrape website content if URL provided
        website_context = ""
        if client_website_url:
            with st.spinner(f"Scraping website: {client_website_url}..."):
                website_context, status = scrape_website_content(client_website_url)
            
            if status == "success":
                st.success(f"✓ Scraped {len(website_context):,} characters from client website.")
            elif status.startswith("error"):
                st.warning(f"⚠️ Website scraping failed: {status.replace('error: ', '')}. Continuing without website context.")
                website_context = ""

        # Extract text from uploaded files
        source_materials = ""
        if uploaded_files:
            with st.spinner(f"Extracting text from {len(uploaded_files)} file(s)..."):
                source_materials = extract_text_from_uploaded_files(uploaded_files)
            if source_materials and not source_materials.startswith("[Error"):
                st.success(f"✓ Extracted {len(source_materials):,} characters from source materials.")
            elif source_materials.startswith("[Error"):
                st.warning("Some files could not be processed. Continuing with available data.")

        with st.spinner("Generating guide with OpenAI..."):
            guide_text = generate_guide(
                product=product,
                goal=goal,
                target_group=target_group,
                personas=personas,
                tone_of_voice=tone_of_voice,
                additional_reference=additional_reference,
                feedback=feedback,
                source_materials=source_materials,
                website_context=website_context,
                sales_methodology=sales_methodology,
                do_not_talk_about=do_not_talk_about,
                strict_qualifying_questions=strict_qualifying_questions,
            )
        
        # Store the generated guide in session state
        st.session_state.last_generated_guide = {
            "guide_text": guide_text,
            "product": product,
            "goal": goal,
            "target_group": target_group,
            "website_context": website_context,
            "client_website_url": client_website_url,
        }
        
        # Display the guide
        st.markdown(guide_text or "")
        
        # Save to Library section
        if guide_text:
            st.divider()
            col1, col2 = st.columns([3, 1])
            with col1:
                guide_name = st.text_input(
                    "Guide Name (optional)",
                    placeholder=f"E.g., '{product[:30]}...' or 'Campaign Q1 2026'",
                    help="Give this guide a memorable name before saving to library."
                )
            with col2:
                st.write("")
                st.write("")
                if st.button("💾 Save to Library", use_container_width=True):
                    from datetime import datetime
                    
                    saved_guide = {
                        "id": len(st.session_state.saved_guides) + 1,
                        "name": guide_name.strip() if guide_name.strip() else f"Guide #{len(st.session_state.saved_guides) + 1}",
                        "date": datetime.now().strftime("%Y-%m-%d %H:%M"),
                        "guide_text": guide_text,
                        "product": product,
                        "goal": goal,
                        "target_group": target_group,
                        "personas": personas,
                    }
                    st.session_state.saved_guides.append(saved_guide)
                    st.success(f"✅ Saved '{saved_guide['name']}' to library!")
                    st.balloons()
        
        # Word Export
        if guide_text:
            # Create clean client info summary for Word template
            client_info_summary = ""
            if website_context and website_context.strip():
                with st.spinner("Summarizing website content..."):
                    # Use AI to create a clean, concise summary
                    website_summary = summarize_website_content(website_context, client_website_url)
                
                if website_summary:
                    client_info_summary = f"Website: {client_website_url}\n\n{website_summary}"
            
            # Always attempt export
            try:
                with st.spinner("Preparing Word export..."):
                    docx_bytes, export_status = export_to_word(guide_text, client_info_summary)
                
                # Show download button if we got valid bytes (even with warnings)
                if docx_bytes and len(docx_bytes) > 0:
                    # Display status message
                    if export_status.startswith("warning"):
                        st.warning(f"⚠️ {export_status}")
                    elif export_status == "success":
                        st.success("✅ Mall laddad och formatering fixad! Dokument redo för nedladdning.")
                    
                    # ALWAYS show download button if we have valid bytes
                    st.download_button(
                        label="📄 Download as Word (.docx)",
                        data=docx_bytes,
                        file_name="discussion_guide.docx",
                        mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                        use_container_width=True,
                    )
                else:
                    # Only show error if truly failed
                    error_msg = export_status.replace("error: ", "") if export_status.startswith("error") else "Unknown error"
                    st.error(f"❌ Word export failed: {error_msg}")
            except Exception as e:
                st.error(f"❌ Unexpected error during Word export: {str(e)}")
        
        with st.expander("Copyable text"):
            st.text_area("Script (raw)", guide_text, height=320)
    else:
        st.info("Fill out the inputs and click Generate Script to generate a cold calling script.")


def render_library_view() -> None:
    """Render the My Library view showing saved guides."""
    st.title("📚 My Library")
    st.caption("Browse and manage your saved discussion guides.")
    
    if not st.session_state.saved_guides:
        st.info("📭 No saved guides yet. Generate and save a script from the Generator to see it here.")
        if st.button("← Go to Generator"):
            st.session_state.current_view = "Generator"
            st.rerun()
    else:
        st.write(f"**Total saved guides:** {len(st.session_state.saved_guides)}")
        st.divider()
        
        for guide in reversed(st.session_state.saved_guides):
            with st.expander(f"📄 {guide['name']} — {guide['date']}", expanded=False):
                # Guide metadata
                col1, col2, col3 = st.columns(3)
                with col1:
                    st.write(f"**Product:** {guide.get('product', 'N/A')[:100]}...")
                with col2:
                    st.write(f"**Goal:** {guide.get('goal', 'N/A')}")
                with col3:
                    st.write(f"**Target Group:** {guide.get('target_group', 'N/A')[:50]}...")
                
                st.divider()
                
                # Guide content
                st.markdown(guide["guide_text"])
                
                st.divider()
                
                # Actions
                col1, col2 = st.columns([3, 1])
                with col1:
                    st.text_area(
                        "Copyable text",
                        guide["guide_text"],
                        height=150,
                        key=f"guide_text_{guide['id']}"
                    )
                with col2:
                    st.write("")
                    st.write("")
                    if st.button(f"🗑️ Delete", key=f"delete_{guide['id']}", use_container_width=True):
                        st.session_state.saved_guides = [
                            g for g in st.session_state.saved_guides if g["id"] != guide["id"]
                        ]
                        st.success(f"Deleted '{guide['name']}'")
                        st.rerun()


def main() -> None:
    st.set_page_config(page_title="Discussion Guide Generator", layout="wide")
    
    # Initialize session state for navigation and saved guides
    if "current_view" not in st.session_state:
        st.session_state.current_view = "Generator"
    if "saved_guides" not in st.session_state:
        st.session_state.saved_guides = []
    if "last_generated_guide" not in st.session_state:
        st.session_state.last_generated_guide = None
    
    # Initialize input field states to persist across navigation
    if "input_product" not in st.session_state:
        st.session_state.input_product = ""
    if "input_goal" not in st.session_state:
        st.session_state.input_goal = "Leads"
    if "input_target_group" not in st.session_state:
        st.session_state.input_target_group = ""
    if "input_personas" not in st.session_state:
        st.session_state.input_personas = ""
    if "input_strict_qualifying_questions" not in st.session_state:
        st.session_state.input_strict_qualifying_questions = ""
    if "input_client_website_url" not in st.session_state:
        st.session_state.input_client_website_url = ""
    if "input_sales_methodology" not in st.session_state:
        st.session_state.input_sales_methodology = "Standard"
    if "input_tone_of_voice" not in st.session_state:
        st.session_state.input_tone_of_voice = "Professional, conversational, confident; concise and value-led."
    if "input_do_not_talk_about" not in st.session_state:
        st.session_state.input_do_not_talk_about = ""
    if "input_additional_reference" not in st.session_state:
        st.session_state.input_additional_reference = ""
    if "input_feedback" not in st.session_state:
        st.session_state.input_feedback = ""
    
    # Custom CSS for wider sidebar and better styling
    st.markdown("""
        <style>
        /* Wider sidebar */
        [data-testid="stSidebar"][aria-expanded="true"] {
            min-width: 500px;
            max-width: 500px;
        }
        
        /* Sidebar background and styling */
        [data-testid="stSidebar"] {
            background-color: #f8f9fa;
        }
        
        /* Bold labels */
        [data-testid="stSidebar"] label {
            font-weight: 600;
            font-size: 0.95rem;
            color: #1f2937;
        }
        
        /* Better spacing for form elements */
        [data-testid="stSidebar"] .stTextArea,
        [data-testid="stSidebar"] .stTextInput,
        [data-testid="stSidebar"] .stSelectbox {
            margin-bottom: 1rem;
        }
        
        /* Expander styling */
        [data-testid="stSidebar"] .streamlit-expanderHeader {
            background-color: #e5e7eb;
            border-radius: 4px;
            font-weight: 600;
            padding: 0.5rem;
        }
        
        /* Generate button styling */
        [data-testid="stSidebar"] button[kind="primary"],
        [data-testid="stSidebar"] button[kind="secondary"] {
            background-color: #2563eb;
            color: white;
            font-weight: 600;
            border-radius: 6px;
            padding: 0.6rem 1rem;
            border: none;
            box-shadow: 0 2px 4px rgba(0,0,0,0.1);
            transition: all 0.2s;
        }
        
        [data-testid="stSidebar"] button[kind="primary"]:hover,
        [data-testid="stSidebar"] button[kind="secondary"]:hover {
            background-color: #1d4ed8;
            box-shadow: 0 4px 6px rgba(0,0,0,0.15);
            transform: translateY(-1px);
        }
        
        /* Divider styling */
        [data-testid="stSidebar"] hr {
            margin: 1.5rem 0;
            border-color: #d1d5db;
        }
        
        /* Main content area */
        .main .block-container {
            padding-top: 2rem;
            max-width: 1200px;
        }
        
        /* Header styling */
        h1 {
            color: #1f2937;
            font-weight: 700;
        }
        
        /* Subheader in main */
        h2 {
            color: #374151;
            margin-top: 2rem;
        }
        
        /* Navigation tabs styling */
        .nav-tabs {
            display: flex;
            justify-content: center;
            gap: 1rem;
            margin-bottom: 2rem;
            padding: 1rem 0;
            border-bottom: 2px solid #e5e7eb;
        }
        
        .nav-tab {
            padding: 0.5rem 2rem;
            background-color: transparent;
            border: none;
            border-bottom: 3px solid transparent;
            font-size: 1.1rem;
            font-weight: 600;
            color: #6b7280;
            cursor: pointer;
            transition: all 0.2s;
        }
        
        .nav-tab.active {
            color: #2563eb;
            border-bottom-color: #2563eb;
        }
        
        .nav-tab:hover {
            color: #1d4ed8;
        }
        </style>
    """, unsafe_allow_html=True)

    # Top navigation bar
    col1, col2, col3 = st.columns([1, 2, 1])
    with col2:
        nav_col1, nav_col2 = st.columns(2)
        with nav_col1:
            if st.button("🎨 Generator", use_container_width=True, type="primary" if st.session_state.current_view == "Generator" else "secondary"):
                st.session_state.current_view = "Generator"
                st.rerun()
        with nav_col2:
            if st.button("📚 My Library", use_container_width=True, type="primary" if st.session_state.current_view == "My Library" else "secondary"):
                st.session_state.current_view = "My Library"
                st.rerun()
    
    st.divider()
    
    # Display the appropriate view
    if st.session_state.current_view == "Generator":
        render_generator_view()
    else:
        render_library_view()


if __name__ == "__main__":
    main()

